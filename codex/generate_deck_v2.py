#!/usr/bin/env python3
"""Generate Codex Enterprise Enablement Slide Deck v2 — VISIBILITY thesis.

7 slides implementing the VISIBILITY thesis for Codex enterprise enablement.
Sales enablement presentation for a Fortune 500 retail deployment.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# -- Color Palette (professional, OpenAI-aligned) --
DARK_BG = RGBColor(0x0D, 0x0D, 0x0D)       # Near-black
DARK_NAVY = RGBColor(0x1A, 0x1A, 0x2E)
OAI_GREEN = RGBColor(0x10, 0xA3, 0x7F)      # OpenAI brand green
ACCENT_BLUE = RGBColor(0x3A, 0x6B, 0x9F)
ACCENT_ORANGE = RGBColor(0xE7, 0x6F, 0x51)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
ACCENT_PURPLE = RGBColor(0x6C, 0x5C, 0xE7)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MED_GRAY = RGBColor(0x88, 0x88, 0x88)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_WHITE = RGBColor(0xF8, 0xF9, 0xFA)
CALLOUT_BG = RGBColor(0xE8, 0xF5, 0xF0)     # Light green tint
VP_CALLOUT = RGBColor(0xFD, 0xF2, 0xE9)      # Light orange tint
DIR_CALLOUT = RGBColor(0xE8, 0xEE, 0xF8)     # Light blue tint


# ================================================================
# UTILITY FUNCTIONS
# ================================================================

def add_bg_rect(slide, color=DARK_BG):
    """Add a full-slide background rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, italic=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri'):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, font_size=16,
                   color=DARK_TEXT, line_spacing=1.3, font_name='Calibri',
                   alignment=PP_ALIGN.LEFT):
    """Add text box with multiple paragraphs. lines = [(text, bold, italic), ...]"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bld, ital) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bld
        p.font.italic = ital
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * (line_spacing - 1) + 4)
    return txBox


def add_accent_bar(slide, left, top, width, height, color=OAI_GREEN):
    """Add a thin accent bar."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_box(slide, left, top, width, height, border_color=OAI_GREEN, bg_color=WHITE):
    """Add a rounded rectangle box."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    return shape


def stat_box(slide, left, top, width, height, number, label,
             num_color=OAI_GREEN, bg_color=WHITE):
    """Add a stat callout box."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = num_color
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(11)
    p2.font.color.rgb = MED_GRAY
    p2.font.name = 'Calibri'
    p2.alignment = PP_ALIGN.CENTER
    return shape


def add_circle(slide, left, top, size, fill_color):
    """Add a filled circle (oval)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(left), Inches(top),
        Inches(size), Inches(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


# ================================================================
# SLIDE 1: TITLE — "Engineering Visibility"
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, DARK_BG)

# Top accent bar
add_accent_bar(slide, 0, 0, 13.333, 0.05, OAI_GREEN)

# Title
add_text_box(slide, 1.5, 1.6, 10, 1.0,
             'Engineering Visibility',
             font_size=48, bold=True, color=WHITE)
add_accent_bar(slide, 1.5, 2.8, 3.0, 0.06, OAI_GREEN)

# Subtitle
add_text_box(slide, 1.5, 3.1, 10, 0.8,
             'How Codex Makes Your Codebase \u2014 and Your Engineering \u2014 Visible',
             font_size=22, color=RGBColor(0xCC, 0xCC, 0xDD))

# Session type
add_text_box(slide, 1.5, 4.0, 10, 0.6,
             'Enterprise Enablement Session',
             font_size=16, italic=True, color=RGBColor(0x99, 0x99, 0xBB))

# Dual audience callout
add_text_box(slide, 1.5, 5.2, 5, 0.35,
             'Prepared for:', font_size=12, bold=True, color=OAI_GREEN)
add_text_box(slide, 1.5, 5.5, 5, 0.35,
             'Director of Engineering  |  VP of IT',
             font_size=14, color=RGBColor(0xBB, 0xBB, 0xCC))
add_text_box(slide, 1.5, 5.9, 5, 0.35,
             'March 2026  |  OpenAI Enterprise',
             font_size=12, color=RGBColor(0x77, 0x77, 0x99))


# ================================================================
# SLIDE 2: WHAT WE HEARD + THE INSIGHT
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, SOFT_WHITE)
add_accent_bar(slide, 0, 0, 13.333, 0.07, OAI_GREEN)
add_text_box(slide, 0.8, 0.4, 11, 0.8, 'Four Problems. One Root Cause.',
             font_size=28, bold=True, color=DARK_NAVY)
add_accent_bar(slide, 0.8, 1.15, 2.5, 0.04, OAI_GREEN)

# Left column — What We Heard
add_text_box(slide, 0.8, 1.5, 5.5, 0.5, 'What We Heard in Discovery',
             font_size=20, bold=True, color=ACCENT_BLUE)
add_multi_text(slide, 0.8, 2.1, 5.5, 3.8, [
    ('\u2022  15-year-old checkout code nobody understands', False, False),
    ('', False, False),
    ('\u2022  Test coverage gaps \u2014 teams afraid to refactor', False, False),
    ('', False, False),
    ('\u2022  Prior AI tool rollout ended in a security incident', False, False),
    ('', False, False),
    ('\u2022  No visibility into what developers were generating', False, False),
], font_size=15, color=DARK_TEXT, line_spacing=1.1)

# Divider
add_accent_bar(slide, 6.45, 1.5, 0.03, 4.5, MED_GRAY)

# Right column — The Insight
add_text_box(slide, 6.8, 1.5, 5.8, 0.5, 'The Insight: One Root Cause',
             font_size=20, bold=True, color=ACCENT_ORANGE)
add_multi_text(slide, 6.8, 2.1, 5.8, 3.8, [
    ('\u2192  Invisible codebase \u2192 devs can\'t understand it', False, False),
    ('', False, False),
    ('\u2192  Invisible quality \u2192 teams can\'t prioritize', False, False),
    ('', False, False),
    ('\u2192  Invisible usage \u2192 leadership can\'t govern', False, False),
    ('', False, False),
    ('\u2192  Invisible impact \u2192 business can\'t justify', False, False),
], font_size=15, color=DARK_TEXT, line_spacing=1.1)

# Bottom callout bar
add_box(slide, 0.5, 6.2, 12.3, 0.9, OAI_GREEN, CALLOUT_BG)
add_text_box(slide, 0.7, 6.25, 12.0, 0.8,
             'All four problems share one root cause: invisibility. '
             'Codex doesn\'t just write code \u2014 it generates the data that makes engineering visible.',
             font_size=15, bold=True, color=DARK_NAVY, alignment=PP_ALIGN.CENTER)


# ================================================================
# SLIDE 3: HOW CODEX WORKS — AND WHY IT'S DIFFERENT
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, SOFT_WHITE)
add_accent_bar(slide, 0, 0, 13.333, 0.07, OAI_GREEN)
add_text_box(slide, 0.8, 0.4, 11, 0.8, 'Agentic. Sandboxed. Every Action Audited.',
             font_size=28, bold=True, color=DARK_NAVY)
add_accent_bar(slide, 0.8, 1.15, 2.5, 0.04, OAI_GREEN)

# Left box — DEVELOPER MACHINE
add_box(slide, 0.5, 1.6, 5.5, 4.0, ACCENT_BLUE, DIR_CALLOUT)
add_text_box(slide, 0.7, 1.7, 5.1, 0.4, 'DEVELOPER MACHINE',
             font_size=14, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# Code Repository sub-box
add_box(slide, 0.8, 2.3, 4.9, 0.9, MED_GRAY, WHITE)
add_text_box(slide, 0.95, 2.35, 4.6, 0.35, 'Code Repository',
             font_size=13, bold=True, color=DARK_TEXT)
add_text_box(slide, 0.95, 2.7, 4.6, 0.45, 'Source files, dependencies, tests, AGENTS.md',
             font_size=11, color=MED_GRAY)

# Codex CLI sub-box
add_box(slide, 0.8, 3.4, 4.9, 0.9, OAI_GREEN, WHITE)
add_text_box(slide, 0.95, 3.45, 4.6, 0.35, 'Codex CLI (Open-Source Rust)',
             font_size=13, bold=True, color=OAI_GREEN)
add_text_box(slide, 0.95, 3.8, 4.6, 0.45, 'Read \u2192 Plan \u2192 Edit \u2192 Run \u2192 Verify',
             font_size=11, color=MED_GRAY)

# Sandbox sub-box
add_box(slide, 0.8, 4.5, 4.9, 0.9, ACCENT_ORANGE, VP_CALLOUT)
add_text_box(slide, 0.95, 4.55, 4.6, 0.35, 'Kernel-Level Sandbox',
             font_size=13, bold=True, color=ACCENT_ORANGE)
add_text_box(slide, 0.95, 4.9, 4.6, 0.45, 'Seatbelt (macOS) / Landlock + seccomp (Linux)',
             font_size=11, color=DARK_TEXT)

# Arrow
add_text_box(slide, 6.05, 2.8, 1.2, 0.6, '\u2192',
             font_size=42, bold=True, color=OAI_GREEN, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 5.95, 3.45, 1.4, 0.5, 'Relevant\ncontext only',
             font_size=10, italic=True, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Right box — OPENAI CLOUD
add_box(slide, 7.3, 1.6, 5.5, 4.0, OAI_GREEN, CALLOUT_BG)
add_text_box(slide, 7.5, 1.7, 5.1, 0.4, 'OPENAI CLOUD',
             font_size=14, bold=True, color=OAI_GREEN, alignment=PP_ALIGN.CENTER)

# GPT-5.4 sub-box
add_box(slide, 7.6, 2.3, 4.9, 0.9, OAI_GREEN, WHITE)
add_text_box(slide, 7.75, 2.35, 4.6, 0.35, 'GPT-5.4 (1M context)',
             font_size=13, bold=True, color=OAI_GREEN)
add_text_box(slide, 7.75, 2.7, 4.6, 0.45, 'Reasoning + code generation. No training on your data.',
             font_size=11, color=MED_GRAY)

# Enterprise Controls sub-box
add_box(slide, 7.6, 3.4, 4.9, 2.0, ACCENT_PURPLE, WHITE)
add_text_box(slide, 7.75, 3.45, 4.6, 0.35, 'Enterprise Controls',
             font_size=13, bold=True, color=ACCENT_PURPLE)
add_text_box(slide, 7.75, 3.85, 4.6, 1.4,
             '\u2022  SAML SSO + SCIM provisioning\n'
             '\u2022  RBAC + Compliance API + audit trails\n'
             '\u2022  Usage analytics dashboard\n'
             '\u2022  requirements.toml \u2014 cloud-managed policies',
             font_size=11, color=DARK_TEXT)

# Key differentiator callout
add_box(slide, 0.5, 5.85, 12.3, 0.7, OAI_GREEN, CALLOUT_BG)
add_text_box(slide, 0.7, 5.9, 12.0, 0.6,
             'Codex is governed AND visible. Open-source CLI, auditable, Compliance API. '
             'Part of OpenAI\'s enterprise suite alongside ChatGPT Enterprise + API + Frontier.',
             font_size=14, bold=True, color=DARK_NAVY, alignment=PP_ALIGN.CENTER)

# VP assurance bar at bottom
add_box(slide, 0.5, 6.7, 12.3, 0.55, ACCENT_ORANGE, VP_CALLOUT)
add_text_box(slide, 0.7, 6.72, 12.0, 0.5,
             'VP Assurance: Source code stays on the developer machine. '
             'Only relevant context reaches the API. Kernel-level sandbox blocks unauthorized access.',
             font_size=12, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)


# ================================================================
# SLIDE 4: THE HARNESS — GOVERNANCE ARCHITECTURE
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, SOFT_WHITE)
add_accent_bar(slide, 0, 0, 13.333, 0.07, ACCENT_PURPLE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, 'Your Standards Become the AI\u2019s Operating Rules',
             font_size=28, bold=True, color=DARK_NAVY)
add_accent_bar(slide, 0.8, 1.15, 2.5, 0.04, ACCENT_PURPLE)

# Three column boxes
harness_items = [
    ('AGENTS.md',
     'Per-repo config, version-controlled',
     'Team-owned',
     ['\u2022  Define what Codex can/cannot do per repo',
      '\u2022  Cascading: org \u2192 repo \u2192 folder',
      '\u2022  Version-controlled alongside code',
      '\u2022  Engineering teams own their policies',
      '',
      'Example:',
      '  "Do not modify files in /payments"',
      '  "Always run tests before committing"'],
     OAI_GREEN),
    ('requirements.toml',
     'Cloud-managed guardrails',
     'IT-owned',
     ['\u2022  Set by IT/Security, enforced centrally',
      '\u2022  Overrides local AGENTS.md when stricter',
      '\u2022  Define allowed/blocked operations globally',
      '\u2022  Audit-friendly: policy = code',
      '',
      'Example:',
      '  "Block all direct database writes"',
      '  "Require review for security paths"'],
     ACCENT_PURPLE),
    ('Approval Modes',
     'Graduated trust for operations',
     'Phased rollout',
     ['\u2022  on-request \u2014 asks before every action',
      '\u2022  untrusted \u2014 auto-approves safe ops',
      '\u2022  team-by-team \u2014 earned autonomy',
      '',
      'Phase mapping:',
      '  Pilot: on-request',
      '  Expand: untrusted',
      '  Scale: team-by-team decision',
      '',
      'Trust is earned, not assumed.'],
     ACCENT_BLUE),
]

for i, (title, subtitle, owner, items, color) in enumerate(harness_items):
    x = 0.5 + i * 4.2
    add_box(slide, x, 1.5, 3.9, 4.8, color, WHITE)
    add_text_box(slide, x + 0.15, 1.6, 3.6, 0.4, title,
                 font_size=18, bold=True, color=color, font_name='Consolas')
    add_text_box(slide, x + 0.15, 2.05, 3.6, 0.3, subtitle,
                 font_size=12, italic=True, color=MED_GRAY)
    add_text_box(slide, x + 0.15, 2.3, 3.6, 0.25, owner,
                 font_size=11, bold=True, color=color)
    add_accent_bar(slide, x + 0.15, 2.55, 2.0, 0.03, color)
    y_pos = 2.7
    for item in items:
        add_text_box(slide, x + 0.15, y_pos, 3.6, 0.3, item,
                     font_size=11, color=DARK_TEXT)
        y_pos += 0.28

# Credibility callout at bottom
add_box(slide, 0.5, 6.5, 12.3, 0.7, OAI_GREEN, CALLOUT_BG)
add_text_box(slide, 0.7, 6.55, 12.0, 0.6,
             'OpenAI built 1 million lines of production code using this exact harness in 5 months. '
             'The open-source CLI is auditable.',
             font_size=14, bold=True, color=DARK_NAVY, alignment=PP_ALIGN.CENTER)


# ================================================================
# SLIDE 5: ENABLEMENT — THREE WORKFLOWS, PHASED ROLLOUT
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, SOFT_WHITE)
add_accent_bar(slide, 0, 0, 13.333, 0.07, OAI_GREEN)
add_text_box(slide, 0.8, 0.4, 11, 0.8, 'Start Read-Only. Scale on Evidence.',
             font_size=28, bold=True, color=DARK_NAVY)
add_accent_bar(slide, 0.8, 1.15, 2.5, 0.04, OAI_GREEN)

# Three workflow boxes
workflows = [
    ('1. Code Understanding',
     'HIGHLIGHTED \u2014 Start here',
     ['\u2022  "Explain this function and its deps"',
      '\u2022  "What would break if I changed this?"',
      '\u2022  Impact analysis + dependency mapping',
      '',
      'Why first:',
      '  \u2713  Read-only = zero generation risk',
      '  \u2713  Addresses pain #1 (legacy code)',
      '  \u2713  Fastest "aha moment"',
      '  \u2713  Gateway to other workflows',
      '',
      'Cisco: 1,500 hrs/month saved',
      'Onboarding: 91\u219249 days (DX Research)'],
     OAI_GREEN, CALLOUT_BG),
    ('2. Test Generation',
     'Measurable coverage lift',
     ['\u2022  "Generate unit tests for this function"',
      '\u2022  "Add edge cases for the payment flow"',
      '\u2022  Coverage reports before/after',
      '',
      'Why second:',
      '  \u2713  Measurable output (coverage %)',
      '  \u2713  Builds on code understanding',
      '  \u2713  Addresses test gap pain point',
      '  \u2713  Generated tests require review'],
     ACCENT_BLUE, WHITE),
    ('3. Docs & Refactoring',
     'Velocity improvement',
     ['\u2022  "Document this module\'s public API"',
      '\u2022  "Refactor to extract shared logic"',
      '\u2022  Inline documentation generation',
      '',
      'Why third:',
      '  \u2713  Highest autonomy, highest review',
      '  \u2713  Devs have trust by this point',
      '  \u2713  Direct velocity impact',
      '  \u2713  Requires mature harness config'],
     ACCENT_PURPLE, WHITE),
]

for i, (title, subtitle, items, color, bg) in enumerate(workflows):
    x = 0.5 + i * 4.2
    add_box(slide, x, 1.5, 3.9, 4.2, color, bg)
    add_text_box(slide, x + 0.15, 1.6, 3.6, 0.4, title,
                 font_size=16, bold=True, color=color)
    add_text_box(slide, x + 0.15, 2.0, 3.6, 0.3, subtitle,
                 font_size=11, bold=True, italic=True,
                 color=OAI_GREEN if i == 0 else MED_GRAY)
    add_accent_bar(slide, x + 0.15, 2.3, 2.0, 0.03, color)
    y_pos = 2.45
    for item in items:
        add_text_box(slide, x + 0.15, y_pos, 3.6, 0.3, item,
                     font_size=11, color=DARK_TEXT)
        y_pos += 0.26

# Phase timeline
add_accent_bar(slide, 0.5, 5.85, 12.3, 0.03, MED_GRAY)
add_text_box(slide, 0.5, 5.95, 4.0, 0.3, 'WEEKS 1\u20132: Pilot',
             font_size=13, bold=True, color=OAI_GREEN)
add_text_box(slide, 0.5, 6.22, 4.0, 0.35,
             '20\u201330 devs  |  Code Understanding  |  on-request  |  weekly office hours',
             font_size=11, color=DARK_TEXT)
add_text_box(slide, 4.7, 5.95, 4.0, 0.3, 'WEEKS 3\u20136: Expand',
             font_size=13, bold=True, color=ACCENT_BLUE)
add_text_box(slide, 4.7, 6.22, 4.0, 0.35,
             '3 teams  |  + Test Generation  |  untrusted',
             font_size=11, color=DARK_TEXT)
add_text_box(slide, 8.9, 5.95, 4.0, 0.3, 'WEEKS 7\u201312: Scale',
             font_size=13, bold=True, color=ACCENT_PURPLE)
add_text_box(slide, 8.9, 6.22, 4.0, 0.35,
             'All engineering  |  + Docs & Refactoring  |  team-by-team',
             font_size=11, color=DARK_TEXT)

# Limitations callout
add_box(slide, 0.5, 6.6, 7.5, 0.6, ACCENT_ORANGE, VP_CALLOUT)
add_text_box(slide, 0.65, 6.63, 7.2, 0.5,
             'Codex has real limitations. We show you how to work WITH them, not around them.',
             font_size=12, bold=True, color=DARK_TEXT)

# Review expectation
add_box(slide, 8.2, 6.6, 4.6, 0.6, ACCENT_RED, WHITE)
add_text_box(slide, 8.35, 6.63, 4.3, 0.5,
             'All generated code = standard code review. No exceptions.',
             font_size=12, bold=True, color=ACCENT_RED)


# ================================================================
# SLIDE 6: THE VISIBILITY LAYER
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, SOFT_WHITE)
add_accent_bar(slide, 0, 0, 13.333, 0.07, OAI_GREEN)
add_text_box(slide, 0.8, 0.4, 11, 0.8,
             'Same Data. Two Audiences.',
             font_size=28, bold=True, color=DARK_NAVY)
add_accent_bar(slide, 0.8, 1.15, 2.5, 0.04, OAI_GREEN)

# --- Conceptual bubble chart / dashboard mockup ---
# Background panel for the dashboard mockup
add_box(slide, 0.5, 1.5, 7.5, 4.2, MED_GRAY, WHITE)
add_text_box(slide, 0.7, 1.55, 7.1, 0.35, 'VISIBILITY DASHBOARD',
             font_size=11, bold=True, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Axis labels for conceptual chart
add_text_box(slide, 0.6, 5.25, 2.0, 0.3, 'Teams \u2192',
             font_size=10, italic=True, color=MED_GRAY)
add_text_box(slide, 0.55, 1.95, 0.7, 2.0, 'Workflows',
             font_size=10, italic=True, color=MED_GRAY)

# Bubble chart concept — circles of different sizes/colors
# Row labels (workflows)
add_text_box(slide, 0.7, 2.3, 1.5, 0.3, 'Understanding',
             font_size=9, color=MED_GRAY)
add_text_box(slide, 0.7, 3.3, 1.5, 0.3, 'Test Gen',
             font_size=9, color=MED_GRAY)
add_text_box(slide, 0.7, 4.3, 1.5, 0.3, 'Docs/Refactor',
             font_size=9, color=MED_GRAY)

# Column labels (teams)
add_text_box(slide, 2.3, 1.95, 1.2, 0.3, 'Platform',
             font_size=9, bold=True, color=ACCENT_BLUE)
add_text_box(slide, 3.6, 1.95, 1.2, 0.3, 'Checkout',
             font_size=9, bold=True, color=ACCENT_PURPLE)
add_text_box(slide, 4.9, 1.95, 1.2, 0.3, 'Inventory',
             font_size=9, bold=True, color=ACCENT_ORANGE)
add_text_box(slide, 6.2, 1.95, 1.2, 0.3, 'Payments',
             font_size=9, bold=True, color=OAI_GREEN)

# Bubbles — Understanding row (read-only = safe, mostly green)
add_circle(slide, 2.35, 2.25, 0.85, OAI_GREEN)        # Platform — large, safe
add_circle(slide, 3.7, 2.4, 0.65, OAI_GREEN)          # Checkout — medium, safe
add_circle(slide, 4.95, 2.55, 0.5, OAI_GREEN)         # Inventory — small-med, safe
add_circle(slide, 6.3, 2.35, 0.75, OAI_GREEN)         # Payments — large-med, safe

# Bubbles — Test Gen row (mix of safe and review needed)
add_circle(slide, 2.5, 3.3, 0.55, OAI_GREEN)          # Platform — safe
add_circle(slide, 3.75, 3.25, 0.6, ACCENT_ORANGE)     # Checkout — review needed
add_circle(slide, 5.1, 3.45, 0.35, ACCENT_ORANGE)     # Inventory — review needed
add_circle(slide, 6.4, 3.35, 0.5, OAI_GREEN)          # Payments — safe

# Bubbles — Docs row (higher risk for write operations)
add_circle(slide, 2.55, 4.3, 0.4, ACCENT_ORANGE)      # Platform — review needed
add_circle(slide, 3.85, 4.4, 0.3, ACCENT_RED)         # Checkout — flagged
add_circle(slide, 5.15, 4.45, 0.25, ACCENT_ORANGE)    # Inventory — review needed
add_circle(slide, 6.45, 4.3, 0.45, ACCENT_ORANGE)     # Payments — review needed

# Legend for bubble sizes
add_text_box(slide, 1.0, 4.95, 6.5, 0.35,
             'Bubble size = query volume     Color = risk level (green = safe, orange = review needed, red = flagged)',
             font_size=9, italic=True, color=MED_GRAY)

# Label on chart
add_text_box(slide, 2.7, 5.25, 5.0, 0.3, 'Team \u00d7 Workflow \u00d7 Risk',
             font_size=12, bold=True, color=DARK_NAVY, alignment=PP_ALIGN.CENTER)

# --- Two callout boxes on the right ---
# VP View box
add_box(slide, 8.3, 1.5, 4.5, 2.0, ACCENT_ORANGE, VP_CALLOUT)
add_text_box(slide, 8.45, 1.55, 4.2, 0.35, 'VP View',
             font_size=16, bold=True, color=ACCENT_ORANGE)
add_accent_bar(slide, 8.45, 1.9, 1.5, 0.03, ACCENT_ORANGE)
add_multi_text(slide, 8.45, 2.0, 4.2, 1.4, [
    ('\u2022  Usage patterns across teams', False, False),
    ('\u2022  Risk indicators by workflow', False, False),
    ('\u2022  Compliance trails, audit-ready', False, False),
    ('\u2022  Export for security review', False, False),
], font_size=12, color=DARK_TEXT, line_spacing=1.1)

# Director View box
add_box(slide, 8.3, 3.7, 4.5, 2.0, ACCENT_BLUE, DIR_CALLOUT)
add_text_box(slide, 8.45, 3.75, 4.2, 0.35, 'Director View',
             font_size=16, bold=True, color=ACCENT_BLUE)
add_accent_bar(slide, 8.45, 4.1, 1.5, 0.03, ACCENT_BLUE)
add_multi_text(slide, 8.45, 4.2, 4.2, 1.4, [
    ('\u2022  Adoption rates by team', False, False),
    ('\u2022  Bottlenecks and friction points', False, False),
    ('\u2022  Impact metrics (coverage, velocity)', False, False),
    ('\u2022  Team performance trends', False, False),
], font_size=12, color=DARK_TEXT, line_spacing=1.1)

# Bottom callout — shared language
add_box(slide, 0.5, 5.8, 12.3, 0.7, OAI_GREEN, CALLOUT_BG)
add_text_box(slide, 0.7, 5.85, 12.0, 0.6,
             'Same data, different views. Shared language for what\'s happening and whether it\'s working.',
             font_size=15, bold=True, color=DARK_NAVY, alignment=PP_ALIGN.CENTER)

# Accent line at bottom
add_accent_bar(slide, 0.5, 6.65, 12.3, 0.03, OAI_GREEN)
add_text_box(slide, 0.5, 6.75, 12.3, 0.5,
             'We built this dashboard in 48 hours using Codex.',
             font_size=14, bold=True, italic=True, color=OAI_GREEN,
             alignment=PP_ALIGN.CENTER)


# ================================================================
# SLIDE 7: NEXT STEPS
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, DARK_BG)
add_accent_bar(slide, 0, 0, 13.333, 0.05, OAI_GREEN)

add_text_box(slide, 1.5, 0.8, 10, 0.8, 'Next Step: Pilot in Two Weeks',
             font_size=36, bold=True, color=WHITE)
add_accent_bar(slide, 1.5, 1.65, 2.5, 0.06, OAI_GREEN)

# Three next steps with numbered circles
steps = [
    ('1', 'Select Pilot Team',
     'Platform engineering \u2014 most legacy burden, most to gain. '
     'Their success becomes the proof point for expansion. '
     '20\u201330 developers. Code Understanding only.',
     OAI_GREEN),
    ('2', 'IT/Security Workshop \u2014 Day 1, Not Day 30',
     'Joint session with your security team and OpenAI to configure '
     'requirements.toml, SSO integration, data classification, and audit policies. '
     'Security is a partner, not a gate.',
     ACCENT_BLUE),
    ('3', 'Define Pass/Fail Metrics Before the Pilot Starts',
     'Agree on measurable success criteria in writing before anyone '
     'touches the tool. If we hit the metric, evidence drives expansion. '
     'If we don\'t, you have a fair, data-driven decision.',
     ACCENT_PURPLE),
]

for i, (num, title, desc, color) in enumerate(steps):
    y = 2.2 + i * 1.6
    # Number circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(1.5), Inches(y),
        Inches(0.6), Inches(0.6)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text_box(slide, 1.5, y + 0.05, 0.6, 0.5, num,
                 font_size=24, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, 2.3, y, 9.5, 0.5, title,
                 font_size=20, bold=True, color=WHITE)
    add_text_box(slide, 2.3, y + 0.5, 9.5, 0.9, desc,
                 font_size=14, color=RGBColor(0xBB, 0xBB, 0xCC))

# Closing thesis
add_accent_bar(slide, 1.5, 6.3, 10, 0.03, OAI_GREEN)
add_text_box(slide, 1.5, 6.45, 10, 0.7,
             '"Start visible. Stay visible. Scale on evidence."',
             font_size=16, bold=True, italic=True, color=OAI_GREEN,
             alignment=PP_ALIGN.CENTER)


# ================================================================
# SAVE
# ================================================================
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, 'codex_enablement_v2.pptx')
prs.save(output_path)
print(f'Deck saved to: {output_path}')
print(f'Slides: {len(prs.slides)}')
